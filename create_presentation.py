import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Initialize presentation with 16:9 widescreen layout
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Scheme Definitions
COLOR_DARK_SLATE = RGBColor(15, 23, 42)    # Slate 900
COLOR_MUTED_SLATE = RGBColor(71, 85, 105)  # Slate 600
COLOR_BLUE = RGBColor(2, 132, 199)         # Sky 700
COLOR_LIGHT_BG = RGBColor(248, 250, 252)   # Slate 50
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BORDER = RGBColor(226, 232, 240)     # Slate 200

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, category_text="WINLOG FORENSIC ANALYZER"):
    # Category / Tag at top
    tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(0.4))
    tf_tag = tag_box.text_frame
    tf_tag.word_wrap = True
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = category_text.upper()
    p_tag.font.name = 'Segoe UI'
    p_tag.font.size = Pt(10)
    p_tag.font.bold = True
    p_tag.font.color.rgb = COLOR_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = 'Segoe UI Semibold'
    p_title.font.size = Pt(28)
    p_title.font.color.rgb = COLOR_DARK_SLATE

# ==================== SLIDE 1: TITLE SLIDE ====================
slide_layout = prs.slide_layouts[6] # Blank layout
slide1 = prs.slides.add_slide(slide_layout)
set_slide_background(slide1, COLOR_LIGHT_BG)

# Add a decorative accent bar
accent_bar = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5)
)
accent_bar.fill.solid()
accent_bar.fill.fore_color.rgb = COLOR_BLUE
accent_bar.line.fill.background()

# Title text frame
main_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11), Inches(4.5))
tf1 = main_box.text_frame
tf1.word_wrap = True

# Main Title
p1 = tf1.paragraphs[0]
p1.text = "WinLog Forensic Importer & Analyzer"
p1.font.name = 'Segoe UI'
p1.font.size = Pt(44)
p1.font.bold = True
p1.font.color.rgb = COLOR_DARK_SLATE
p1.space_after = Pt(10)

# Subtitle
p2 = tf1.add_paragraph()
p2.text = "Offline Event Log and Data Analysis Utility for Digital Forensics"
p2.font.name = 'Segoe UI'
p2.font.size = Pt(22)
p2.font.color.rgb = COLOR_BLUE
p2.space_after = Pt(35)

# Analyst / Tool Info
p3 = tf1.add_paragraph()
p3.text = "• Multi-Format Support: EVTX, JSON, XML, CSV, Excel\n• Preserves Evidence Integrity (Read-Only Offline Analysis)\n• Interactive UI Matching Microsoft Event Viewer Metadata Layout"
p3.font.name = 'Segoe UI'
p3.font.size = Pt(14)
p3.font.color.rgb = COLOR_MUTED_SLATE

# ==================== SLIDE 2: WORKFLOW & INTEGRITY ====================
slide2 = prs.slides.add_slide(slide_layout)
set_slide_background(slide2, COLOR_WHITE)
add_header(slide2, "1. Forensic Workflow & Evidence Integrity")

# Create 3 columns for 3 key concepts
col_width = Inches(3.6)
col_gap = Inches(0.4)
top_pos = Inches(2.0)
height_pos = Inches(4.5)

concepts = [
    {
        "title": "Evidence Integrity",
        "desc": "Analysis is performed offline. Instead of touching a live target system directly, the investigator imports exported log files.\n\nThis maintains strict forensic integrity and ensures no system files or attributes are modified (Artifact Preservation)."
    },
    {
        "title": "Multi-Format Parser",
        "desc": "Directly imports log data from multiple sources:\n• Native Windows Event Logs (.evtx)\n• Structured JSON system/app events\n• XML exported event logs\n• Firewalls & network logs (CSV)\n• Analyst spreadsheets (.xlsx/.xls)."
    },
    {
        "title": "Smart Header Mapping",
        "desc": "The import engine dynamically maps column headers from CSV and Excel worksheets (such as Date, Time, Created, EventID, Level) to a unified UI structure.\n\nThis automatically merges diverse log sources into a single chronologically sorted timeline."
    }
]

for idx, concept in enumerate(concepts):
    left_pos = Inches(0.8) + idx * (col_width + col_gap)
    # Background card shape
    card = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, col_width, height_pos
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_LIGHT_BG
    card.line.color.rgb = COLOR_BORDER
    card.line.width = Pt(1.5)
    
    # Text frame on card
    card_tf = card.text_frame
    card_tf.word_wrap = True
    card_tf.margin_left = Inches(0.2)
    card_tf.margin_right = Inches(0.2)
    card_tf.margin_top = Inches(0.2)
    
    # Title
    p_title = card_tf.paragraphs[0]
    p_title.text = concept["title"]
    p_title.font.name = 'Segoe UI Semibold'
    p_title.font.size = Pt(18)
    p_title.font.color.rgb = COLOR_BLUE
    p_title.space_after = Pt(15)
    
    # Description
    p_desc = card_tf.add_paragraph()
    p_desc.text = concept["desc"]
    p_desc.font.name = 'Segoe UI'
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = COLOR_DARK_SLATE

# ==================== SLIDE 3: INTERFACE & METADATA ====================
slide3 = prs.slides.add_slide(slide_layout)
set_slide_background(slide3, COLOR_WHITE)
add_header(slide3, "2. User Interface & Dynamic Quick Filters")

# Two sections: Left (Event Viewer Match), Right (Dynamic Presets)
panel_width = Inches(5.6)
panel_gap = Inches(0.5)

# Panel 1: Event Viewer Layout
p1_left = Inches(0.8)
card1 = slide3.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, p1_left, top_pos, panel_width, height_pos
)
card1.fill.solid()
card1.fill.fore_color.rgb = COLOR_LIGHT_BG
card1.line.color.rgb = COLOR_BORDER
card1.line.width = Pt(1.5)

tf1 = card1.text_frame
tf1.word_wrap = True
tf1.margin_left = Inches(0.3)
tf1.margin_top = Inches(0.3)

p1_title = tf1.paragraphs[0]
p1_title.text = "Event Viewer Metadata Matching"
p1_title.font.name = 'Segoe UI Semibold'
p1_title.font.size = Pt(18)
p1_title.font.color.rgb = COLOR_BLUE
p1_title.space_after = Pt(15)

p1_desc = tf1.add_paragraph()
p1_desc.text = (
    "The details sidebar displays event data in the official Event Viewer metadata format:\n\n"
    "• Top Section: Shows the full multi-line description text (Message).\n\n"
    "• Separator Line: Clean visual divider.\n\n"
    "• Two-Column Metadata Grid:\n"
    "  Organizes key variables just like Windows: Log Name, Source, Event ID, Level, User, Logged (Timestamp), Task Category, Keywords, Computer name, and OpCode."
)
p1_desc.font.name = 'Segoe UI'
p1_desc.font.size = Pt(13.5)
p1_desc.font.color.rgb = COLOR_DARK_SLATE

# Panel 2: Dynamic Presets
p2_left = Inches(0.8) + panel_width + panel_gap
card2 = slide3.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, p2_left, top_pos, panel_width, height_pos
)
card2.fill.solid()
card2.fill.fore_color.rgb = COLOR_LIGHT_BG
card2.line.color.rgb = COLOR_BORDER
card2.line.width = Pt(1.5)

tf2 = card2.text_frame
tf2.word_wrap = True
tf2.margin_left = Inches(0.3)
tf2.margin_top = Inches(0.3)

p2_title = tf2.paragraphs[0]
p2_title.text = "Dynamic Quick Filters"
p2_title.font.name = 'Segoe UI Semibold'
p2_title.font.size = Pt(18)
p2_title.font.color.rgb = COLOR_BLUE
p2_title.space_after = Pt(15)

p2_desc = tf2.add_paragraph()
p2_desc.text = (
    "Replaces static buttons with a live, adaptive scanning system:\n\n"
    "• Adaptive Presets: Forensic buttons (e.g. Log Clear 104, Failed Logons 4625, Process Creation 4688, Startup/Shutdown 6005/6006/1074) are enabled only if their corresponding Event IDs exist in the file. Disabled buttons show tooltips to avoid confusion.\n\n"
    "• Multi-Log Quick Filters:\n"
    "  When opening custom logs (e.g., WLAN, OpenSSH, or CSV/Excel files), the program automatically extracts the Top 5 Event IDs and Top 5 Sources, dynamically creating clickable filter buttons in the sidebar."
)
p2_desc.font.name = 'Segoe UI'
p2_desc.font.size = Pt(13)
p2_desc.font.color.rgb = COLOR_DARK_SLATE


# ==================== SLIDE 4: PORTABILITY & TIMELINE ====================
slide4 = prs.slides.add_slide(slide_layout)
set_slide_background(slide4, COLOR_WHITE)
add_header(slide4, "3. Playbooks, Export & Portability")

# 3 columns for Slide 4
concepts_s4 = [
    {
        "title": "Forensic Advisor",
        "desc": "Integrates forensic guidance directly into the application.\n\nSelecting a critical security event (like failed logon 4625 or service persistence 7045) immediately displays security implications and step-by-step investigation playbooks for the analyst."
    },
    {
        "title": "Timeline Export",
        "desc": "After completing searches, text filtering, and alert isolation, the analyst can export the current view with a single click.\n\nSupports exporting the filtered timeline to CSV or JSON formats, ready for inclusion in official reports."
    },
    {
        "title": "Portable Deployment",
        "desc": "Built in C# WPF to compile into a single-file executable (Single-File Exe).\n\nCan run directly from a USB forensics drive on any workstation without installing .NET frameworks, dependencies, or local databases."
    }
]

for idx, concept in enumerate(concepts_s4):
    left_pos = Inches(0.8) + idx * (col_width + col_gap)
    # Background card shape
    card = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, col_width, height_pos
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_LIGHT_BG
    card.line.color.rgb = COLOR_BORDER
    card.line.width = Pt(1.5)
    
    # Text frame on card
    card_tf = card.text_frame
    card_tf.word_wrap = True
    card_tf.margin_left = Inches(0.2)
    card_tf.margin_right = Inches(0.2)
    card_tf.margin_top = Inches(0.2)
    
    # Title
    p_title = card_tf.paragraphs[0]
    p_title.text = concept["title"]
    p_title.font.name = 'Segoe UI Semibold'
    p_title.font.size = Pt(18)
    p_title.font.color.rgb = COLOR_BLUE
    p_title.space_after = Pt(15)
    
    # Description
    p_desc = card_tf.add_paragraph()
    p_desc.text = concept["desc"]
    p_desc.font.name = 'Segoe UI'
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = COLOR_DARK_SLATE

# Save the presentation
output_path = r"c:\Users\pc\Desktop\WinLog\WinLog_Forensics_Presentation.pptx"
try:
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")
except PermissionError:
    output_path = r"c:\Users\pc\Desktop\WinLog\WinLog_Forensics_Presentation_v2.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")
